from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation("certificate.pptx")
slide = prs.slides[0]

left = Inches(0.98)
top = Inches(2.02)
height = Inches(3.72)
width = Inches(11.37)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

p = tf.add_paragraph()
student_name = ""
student_college = ""
p.text = "This is to certify that " + student_name + "\nof " + student_college + " participated\nin Vocal Solo at the inter-collegiate festival\n of Ashoka University, Banjaara 2017."
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.size = Pt(30)
p.font.name = "Lato"
p.line_spacing = 1.5
p.alignment = PP_ALIGN.CENTER

prs.save('certificate.pptx')









